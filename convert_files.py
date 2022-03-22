## Converts select files into text files to be parsed over

import collections 
import collections.abc
from pptx import Presentation
import docx
from docx import Document




class Convert:

    def __init__(self, file1, file2):
        self.file1 = file1
        self.file2 = file2
      

    def test(self):
        print("test is working")

    def convert_from_pdf(self):
        return 0

    def convert_from_word(self):
       
        # open connection to Word Document
        doc = docx.Document(self.file1)
        
        # read in each paragraph in file
        result = [p.text for p in doc.paragraphs]
        
        for val in result:
            print(val)

        return result
        

    def convert_from_powerpoint(self):
        prs = Presentation(self.file2)
        text_runs = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
        print(text_runs)

    def create_word_doc(self):
        #Will be used to results back into a word document
        document = Document()
        document.add_heading('The REAL meaning of the universe')
        document.add_paragraph('Lorem ipsum dolor sit amet.')

        document.save('./datafiles/results/test.docx')
        
       
